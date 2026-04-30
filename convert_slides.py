#!/usr/bin/env python3
"""
HTML Slides to Editable PPTX Converter

Supports two HTML presentation styles:
  - Q&A style (GPU.html): nav tabs + collapsible theory/task/solution blocks
  - Slide style (GIMP.html, unittest.html, etc.): sequential .slide divs

Usage:
    python convert_slides.py <file_or_directory> [<file2> ...]

Requirements:
    playwright, python-pptx, beautifulsoup4, lxml
"""

import argparse
import os
import time
from pathlib import Path
from typing import List, Dict

from playwright.sync_api import sync_playwright
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from bs4 import BeautifulSoup


# ─────────────────────────────────────────────
# Detection
# ─────────────────────────────────────────────

def detect_style(html_path: str) -> str:
    """Return 'qa' if the file uses Q&A nav-tab style, else 'slide'."""
    with open(html_path, encoding='utf-8', errors='ignore') as f:
        content = f.read()
    if 'data-pres=' in content and 'toggle-btn' in content:
        return 'qa'
    return 'slide'


# ─────────────────────────────────────────────
# Extraction — Q&A style
# ─────────────────────────────────────────────

def extract_qa(html_path: str) -> List[Dict]:
    """Extract content from Q&A style presentations (GPU.html)."""
    slides_data = []

    with sync_playwright() as p:
        browser = p.chromium.launch(headless=True)
        page = browser.new_page()
        page.goto(f"file://{os.path.abspath(html_path)}")
        page.wait_for_load_state('networkidle')

        nav_buttons = page.locator('nav button').all()
        print(f"  Found {len(nav_buttons)} presentations")

        for idx in range(1, len(nav_buttons) + 1):
            page.locator(f'nav button[data-pres="{idx}"]').click()
            time.sleep(0.3)

            # Expand all collapsible blocks
            for btn in page.locator('button.toggle-btn').all():
                try:
                    btn.click()
                    time.sleep(0.05)
                except Exception:
                    pass
            time.sleep(0.3)

            soup = BeautifulSoup(page.content(), 'lxml')
            active = soup.find('div', class_='presentation active')
            if not active:
                continue

            title_el = active.find('h2')
            pres = {
                'presentation_number': idx,
                'title': title_el.get_text(strip=True) if title_el else f'Презентация {idx}',
                'questions': []
            }

            for card in active.find_all('div', class_='card'):
                q: Dict = {}
                header = card.find('div', class_='card-header')
                if header:
                    q['question'] = header.get_text(strip=True)

                for block in card.find_all('div', id=lambda x: x and x.startswith('theory-')):
                    paras = [p.get_text(strip=True) for p in block.find_all('p')]
                    if paras:
                        q['theory'] = '\n'.join(paras)

                for block in card.find_all('div', id=lambda x: x and x.startswith('task-')):
                    paras = [p.get_text(strip=True) for p in block.find_all('p')]
                    if paras:
                        q['task'] = '\n'.join(paras)

                for block in card.find_all('div', id=lambda x: x and x.startswith('sol-')):
                    sols = [s.get_text(strip=True) for s in block.find_all('div', class_='solution')]
                    if sols:
                        q['solution'] = '\n'.join(sols)

                if q:
                    pres['questions'].append(q)

            slides_data.append(pres)
            print(f"  Presentation {idx}: {len(pres['questions'])} questions")

        browser.close()

    return slides_data


# ─────────────────────────────────────────────
# Extraction — Slide style
# ─────────────────────────────────────────────

def extract_slides(html_path: str) -> List[Dict]:
    """Extract content from sequential-slide presentations."""
    with sync_playwright() as p:
        browser = p.chromium.launch(headless=True)
        page = browser.new_page(viewport={'width': 1920, 'height': 1080})
        page.goto(f"file://{os.path.abspath(html_path)}")
        page.wait_for_load_state('networkidle')
        time.sleep(0.5)
        soup = BeautifulSoup(page.content(), 'lxml')
        browser.close()

    slides_data = []
    for idx, slide in enumerate(soup.find_all('div', class_='slide')):
        # strip scripts/styles
        for tag in slide.find_all(['script', 'style']):
            tag.decompose()

        title_el = slide.find(['h1', 'h2']) or slide.find(class_='slide-title')
        title = title_el.get_text(strip=True) if title_el else ''

        content = []

        for h in slide.find_all(['h1', 'h2', 'h3', 'h4']):
            content.append({'type': 'header', 'level': int(h.name[1]), 'text': h.get_text(strip=True)})

        for p in slide.find_all('p'):
            text = p.get_text(strip=True)
            if text:
                content.append({'type': 'paragraph', 'text': text})

        for ul in slide.find_all(['ul', 'ol']):
            items = [li.get_text(strip=True) for li in ul.find_all('li', recursive=False)]
            if items:
                content.append({'type': 'list', 'items': items})

        for pre in slide.find_all(['pre', 'code']):
            text = pre.get_text(strip=True)
            if len(text) > 10:
                content.append({'type': 'code', 'text': text})

        for table in slide.find_all('table'):
            rows = []
            for tr in table.find_all('tr'):
                cells = [td.get_text(strip=True) for td in tr.find_all(['th', 'td'])]
                if cells:
                    rows.append(cells)
            if rows:
                content.append({'type': 'table', 'rows': rows})

        if title or content:
            slides_data.append({'slide_number': idx + 1, 'title': title, 'content': content})
            print(f"  Slide {idx + 1}: {title[:60] or 'Untitled'}")

    return slides_data


# ─────────────────────────────────────────────
# PPTX creation — Q&A style
# ─────────────────────────────────────────────

def build_pptx_qa(slides_data: List[Dict], output_path: str):
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    BLUE   = RGBColor(37,  99,  235)
    AMBER  = RGBColor(245, 158,  11)
    GREEN  = RGBColor(16,  185, 129)
    DARK   = RGBColor(15,   23,  42)

    for pres in slides_data:
        # Title slide
        ts = prs.slides.add_slide(prs.slide_layouts[0])
        ts.shapes.title.text = pres['title']
        ts.shapes.title.text_frame.paragraphs[0].font.size = Pt(32)
        ts.shapes.title.text_frame.paragraphs[0].font.bold = True
        ts.shapes.title.text_frame.paragraphs[0].font.color.rgb = BLUE
        ts.placeholders[1].text = f"Презентация {pres['presentation_number']}"

        for q in pres['questions']:
            slide = prs.slides.add_slide(prs.slide_layouts[5])

            # Question title
            tb = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
            tb.text_frame.word_wrap = True
            p = tb.text_frame.paragraphs[0]
            p.text = q.get('question', '')
            p.font.size = Pt(18)
            p.font.bold = True
            p.font.color.rgb = DARK

            # Body
            body = slide.shapes.add_textbox(Inches(0.5), Inches(1.3), Inches(9), Inches(5.7))
            body.text_frame.word_wrap = True
            tf = body.text_frame

            def add_section(label: str, text: str, color: RGBColor):
                p = tf.add_paragraph()
                p.text = label
                p.font.size = Pt(13)
                p.font.bold = True
                p.font.color.rgb = color
                p.space_after = Pt(4)

                p = tf.add_paragraph()
                p.text = text
                p.font.size = Pt(10)
                p.space_after = Pt(10)

            first = True
            if 'theory' in q:
                p = tf.paragraphs[0] if first else tf.add_paragraph()
                first = False
                p.text = "📖 Теория:"
                p.font.size = Pt(13)
                p.font.bold = True
                p.font.color.rgb = BLUE
                p.space_after = Pt(4)
                p = tf.add_paragraph()
                p.text = q['theory']
                p.font.size = Pt(10)
                p.space_after = Pt(10)

            if 'task' in q:
                add_section("📝 Практическое задание:", q['task'], AMBER)

            if 'solution' in q:
                add_section("✅ Решение:", q['solution'], GREEN)

    prs.save(output_path)


# ─────────────────────────────────────────────
# PPTX creation — Slide style
# ─────────────────────────────────────────────

def build_pptx_slides(slides_data: List[Dict], output_path: str):
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    BLUE = RGBColor(37, 99, 235)
    STEEL = RGBColor(75, 139, 190)

    for info in slides_data:
        slide = prs.slides.add_slide(prs.slide_layouts[5])

        if info['title']:
            tb = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
            tb.text_frame.word_wrap = True
            p = tb.text_frame.paragraphs[0]
            p.text = info['title']
            p.font.size = Pt(26)
            p.font.bold = True
            p.font.color.rgb = BLUE

        body = slide.shapes.add_textbox(Inches(0.5), Inches(1.3), Inches(9), Inches(5.7))
        body.text_frame.word_wrap = True
        tf = body.text_frame
        first = True

        for el in info['content']:
            p = tf.paragraphs[0] if first else tf.add_paragraph()
            first = False

            if el['type'] == 'header':
                p.text = el['text']
                p.font.size = Pt(max(10, 20 - el['level'] * 2))
                p.font.bold = True
                p.font.color.rgb = STEEL
                p.space_after = Pt(6)

            elif el['type'] == 'paragraph':
                p.text = el['text']
                p.font.size = Pt(11)
                p.space_after = Pt(6)

            elif el['type'] == 'list':
                for i, item in enumerate(el['items']):
                    lp = p if (i == 0) else tf.add_paragraph()
                    lp.text = f"• {item}"
                    lp.font.size = Pt(10)
                    lp.space_after = Pt(3)

            elif el['type'] == 'code':
                p.text = el['text'][:600]
                p.font.size = Pt(9)
                p.font.name = 'Courier New'
                p.space_after = Pt(6)

            elif el['type'] == 'table':
                for row in el['rows'][:8]:
                    rp = p if first else tf.add_paragraph()
                    first = False
                    rp.text = ' | '.join(row)
                    rp.font.size = Pt(9)
                    rp.space_after = Pt(3)

    prs.save(output_path)


# ─────────────────────────────────────────────
# Orchestration
# ─────────────────────────────────────────────

def process_file(html_path: Path):
    if not html_path.exists():
        print(f"❌ Not found: {html_path}")
        return

    print(f"\n{'─'*60}")
    print(f"📄 {html_path.name}")
    print(f"{'─'*60}")

    output_dir = Path("output")
    output_dir.mkdir(exist_ok=True)
    out = output_dir / f"{html_path.stem}_editable.pptx"

    style = detect_style(str(html_path))
    print(f"  Style: {style}")

    try:
        if style == 'qa':
            data = extract_qa(str(html_path))
            build_pptx_qa(data, str(out))
        else:
            data = extract_slides(str(html_path))
            build_pptx_slides(data, str(out))

        print(f"  ✅ Saved → {out}")
    except Exception as e:
        print(f"  ❌ Error: {e}")
        import traceback
        traceback.print_exc()


def process_path(path_str: str):
    path = Path(path_str)
    if path.is_dir():
        files = sorted(path.glob("*.html"))
        if not files:
            print(f"No .html files found in {path}")
        for f in files:
            process_file(f)
    else:
        process_file(path)


def main():
    parser = argparse.ArgumentParser(description="Convert HTML presentations to editable PPTX")
    parser.add_argument('paths', nargs='+', help='HTML files or directories')
    for p in parser.parse_args().paths:
        process_path(p)


if __name__ == "__main__":
    main()
